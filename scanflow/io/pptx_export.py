"""Export a SurveyManifest as a PowerPoint deck.

Slide 1 is the wide-area overview (with the annotated PNG produced by the
runner — numbered boxes around each discovered feature). Slides 2..N each
show one feature: its highest-iteration scan preview plus a metadata box
listing position, size, bias, setpoint, and the per-iteration centering
residuals.

ProbeFlow can call :func:`export_pptx` with a custom ``image_resolver`` to
swap in processed images while keeping all metadata intact.
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Callable, Optional

from scanflow.automation.survey import FeatureRecord, SurveyManifest

log = logging.getLogger(__name__)

ImageResolver = Callable[[FeatureRecord], Optional[Path]]


def _default_image_resolver(rec: FeatureRecord) -> Optional[Path]:
    """Use the last iteration's raw preview PNG."""
    if rec.preview_paths:
        return Path(rec.preview_paths[-1])
    return None


def export_pptx(
    manifest: SurveyManifest,
    output_path: Path,
    image_resolver: Optional[ImageResolver] = None,
) -> Path:
    """Build a PPTX from a SurveyManifest and write it to ``output_path``."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    if image_resolver is None:
        image_resolver = _default_image_resolver

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # fully blank

    # --- Slide 1: overview ----------------------------------------------
    slide = prs.slides.add_slide(blank)
    title = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12.5), Inches(0.6))
    tf = title.text_frame
    tf.text = manifest.name
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True

    subtitle = slide.shapes.add_textbox(Inches(0.3), Inches(0.8), Inches(12.5), Inches(0.4))
    subtitle.text_frame.text = (
        f"{manifest.timestamp}   |   "
        f"{manifest.wide_size_nm[0]:.0f} × {manifest.wide_size_nm[1]:.0f} nm wide field   |   "
        f"{len(manifest.features)} feature(s)"
    )
    subtitle.text_frame.paragraphs[0].font.size = Pt(14)

    overview_img = manifest.wide_preview_path or manifest.wide_scan_path
    overview_img = _resolve_annotated(manifest)
    if overview_img and Path(overview_img).exists():
        slide.shapes.add_picture(
            str(overview_img),
            Inches(3.5), Inches(1.4), height=Inches(5.8),
        )
    else:
        _add_missing_image_placeholder(slide, "wide-area overview not available")

    # --- Slide 2..N: per feature ----------------------------------------
    for rec in manifest.features:
        slide = prs.slides.add_slide(blank)
        # Header
        h = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12.5), Inches(0.6))
        h.text_frame.text = f"Feature #{rec.index}    |    size ≈ {rec.char_dim_nm:.2f} nm"
        h.text_frame.paragraphs[0].font.size = Pt(24)
        h.text_frame.paragraphs[0].font.bold = True

        # Image — left half
        img_path = image_resolver(rec)
        if img_path and Path(img_path).exists():
            slide.shapes.add_picture(
                str(img_path),
                Inches(0.4), Inches(1.1), height=Inches(5.8),
            )
        else:
            _add_missing_image_placeholder(slide, f"no preview for feature {rec.index}")

        # Metadata — right half
        meta = slide.shapes.add_textbox(Inches(6.8), Inches(1.2), Inches(6.2), Inches(5.5))
        tf = meta.text_frame
        tf.word_wrap = True
        _add_meta_line(tf, "Position (Δ from centre)", first=True,
                       value=f"{rec.centroid_nm_offset[0]:+.2f} nm, "
                             f"{rec.centroid_nm_offset[1]:+.2f} nm")
        _add_meta_line(tf, "Zoom frame",
                       value=f"{rec.zoom_size_nm[0]:.2f} × {rec.zoom_size_nm[1]:.2f} nm")
        _add_meta_line(tf, "Bias", value=f"{rec.bias_V:.3f} V")
        _add_meta_line(tf, "Setpoint", value=f"{rec.setpoint_A * 1e12:.1f} pA")
        _add_meta_line(tf, "Iterations", value=str(len(rec.scan_paths)))

        for i, (dx, dy) in enumerate(rec.drift_log_angstrom, start=1):
            _add_meta_line(
                tf,
                f"  iter {i} residual",
                value=f"dx = {dx:+.2f} Å,  dy = {dy:+.2f} Å",
                small=True,
            )

        _add_meta_line(
            tf,
            "Final residual",
            value=f"dx = {rec.final_residual_angstrom[0]:+.2f} Å,  "
                  f"dy = {rec.final_residual_angstrom[1]:+.2f} Å",
        )

        # Z stability — show the best (smallest RMS) of the iteration scans
        z_iters = getattr(rec, "z_stability_per_iter", None) or []
        if z_iters:
            best = min(z_iters, key=lambda m: m.get("rms_pm", float("inf")))
            _add_meta_line(
                tf,
                "Z stability",
                value=(
                    f"{best.get('rms_pm', 0):.1f} pm RMS  "
                    f"({best.get('rating', '?')}, "
                    f"{int(best.get('jumps', 0))} jump(s))"
                ),
            )
        if rec.scan_paths:
            _add_meta_line(tf, "Source", value=Path(rec.scan_paths[-1]).name, small=True)

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    log.info("Wrote PPTX: %s", output_path)
    return output_path


def _resolve_annotated(manifest: SurveyManifest) -> Optional[Path]:
    """Prefer the annotated overview (with feature boxes) if available."""
    if manifest.wide_preview_path:
        p = Path(manifest.wide_preview_path)
        annotated = p.with_name(p.stem + "_annotated" + p.suffix)
        # Runner writes wide_annotated.png — check sibling first
        sibling = p.parent / "wide_annotated.png"
        if sibling.exists():
            return sibling
        if annotated.exists():
            return annotated
        if p.exists():
            return p
    return None


def _add_meta_line(tf, label: str, *, value: str = "", first: bool = False,
                   small: bool = False) -> None:
    from pptx.util import Pt
    if first:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = f"{label}: {value}" if value else label
    p.font.size = Pt(10 if small else 13)


def _add_missing_image_placeholder(slide, msg: str) -> None:
    from pptx.util import Inches, Pt
    box = slide.shapes.add_textbox(Inches(3.5), Inches(3.0), Inches(6.0), Inches(1.0))
    box.text_frame.text = f"[ {msg} ]"
    box.text_frame.paragraphs[0].font.size = Pt(16)
